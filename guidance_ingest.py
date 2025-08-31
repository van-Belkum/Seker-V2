import os, io, hashlib, datetime as dt
from typing import Dict, Any, List
import yaml

try:
    import docx
except:
    docx = None
try:
    from pptx import Presentation
except:
    Presentation = None

RULES_DIR = "rules"
PARSED_DIR = os.path.join(RULES_DIR, "parsed")
GUID_DIR = os.path.join("data","guidance")
os.makedirs(PARSED_DIR, exist_ok=True)
os.makedirs(GUID_DIR, exist_ok=True)

def _sha(b: bytes) -> str:
    return hashlib.sha1(b).hexdigest()

def secure_store_guidance(name: str, raw: bytes, visibility: str="private") -> str:
    gid = _sha(raw)
    path = os.path.join(GUID_DIR, f"{gid}.bin")
    with open(path,"wb") as f:
        f.write(raw)
    meta = {"id": gid, "name": name, "stored": dt.datetime.utcnow().isoformat(), "visibility": visibility, "active": True}
    idx = _load_index()
    idx["items"] = [x for x in idx.get("items",[]) if x["id"]!=gid]
    idx["items"].append(meta)
    _save_index(idx)
    return gid

def list_guidance_versions() -> List[Dict[str,Any]]:
    return _load_index().get("items",[])

def supersede_guidance(gid: str) -> bool:
    idx = _load_index()
    ok=False
    for x in idx.get("items",[]):
        if x["id"]==gid:
            x["active"]=False
            ok=True
    _save_index(idx)
    return ok

def _load_index():
    fn = os.path.join(GUID_DIR,"index.yaml")
    if not os.path.exists(fn):
        return {"items":[]}
    return yaml.safe_load(open(fn,"r")) or {"items":[]}

def _save_index(idx):
    fn = os.path.join(GUID_DIR,"index.yaml")
    with open(fn,"w") as f:
        yaml.safe_dump(idx, f, sort_keys=False)

def _extract_text_from_docx(raw: bytes) -> str:
    if not docx:
        return ""
    bio = io.BytesIO(raw)
    d = docx.Document(bio)
    return "\n".join([p.text for p in d.paragraphs])

def _extract_text_from_pptx(raw: bytes) -> str:
    if not Presentation:
        return ""
    bio = io.BytesIO(raw)
    prs = Presentation(bio)
    txt=[]
    for s in prs.slides:
        for shp in s.shapes:
            if hasattr(shp,"text"):
                txt.append(shp.text)
    return "\n".join(txt)

def _extract_text_from_pdf(raw: bytes) -> str:
    try:
        import fitz
        doc = fitz.open(stream=raw, filetype="pdf")
        out=[]
        for p in doc:
            out.append(p.get_text("text") or "")
        doc.close()
        return "\n".join(out)
    except:
        return ""

def parse_guidance_to_rules(name: str, raw: bytes) -> str:
    text=""
    ext = os.path.splitext(name)[1].lower()
    if ext==".docx":
        text=_extract_text_from_docx(raw)
    elif ext==".pptx":
        text=_extract_text_from_pptx(raw)
    elif ext==".pdf":
        text=_extract_text_from_pdf(raw)

    rules=[]
    for line in text.splitlines():
        s=line.strip()
        if not s: continue
        if any(s.lower().startswith(w) for w in ["must","shall","should","ensure","do not","forbidden","require","required"]):
            r = {"title": s[:80], "reason": s, "severity":"major"}
            if "do not" in s.lower() or "forbidden" in s.lower():
                r["must_not_contain"] = s.split("do not")[-1].strip() if "do not" in s.lower() else s
            else:
                r["must_contain"] = s
            rules.append(r)

    if not rules:
        for s in [x.strip() for x in text.split(".") if len(x.strip())>25]:
            rules.append({"title": s[:80], "reason": s, "severity":"minor", "must_contain": s})

    out = {"rules": rules}
    gid = _sha(raw)
    fn = os.path.join(PARSED_DIR, f"{gid}.yaml")
    with open(fn,"w") as f:
        yaml.safe_dump(out, f, sort_keys=False)
    return fn
