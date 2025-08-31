from docx import Document as Docx
from pptx import Presentation
from pypdf import PdfReader

def extract_text_from_docx(path: str) -> str:
    doc = Docx(path)
    return "\n".join(p.text for p in doc.paragraphs)

def extract_text_from_pptx(path: str) -> str:
    prs = Presentation(path)
    txt = []
    for s in prs.slides:
        for shp in s.shapes:
            if hasattr(shp, "text"):
                txt.append(shp.text)
    return "\n".join(txt)

def extract_text_from_pdf(path: str) -> str:
    reader = PdfReader(path)
    chunks = []
    for p in reader.pages:
        try:
            chunks.append(p.extract_text() or "")
        except Exception:
            chunks.append("")
    return "\n".join(chunks)
