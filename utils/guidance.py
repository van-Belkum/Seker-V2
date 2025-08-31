import os
import docx
import pptx
import fitz

def extract_text_from_file(path):
    if path.lower().endswith(".docx"):
        doc = docx.Document(path)
        return "\n".join([p.text for p in doc.paragraphs])
    elif path.lower().endswith(".pptx"):
        prs = pptx.Presentation(path)
        return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    elif path.lower().endswith(".pdf"):
        text = ""
        with fitz.open(path) as doc:
            for page in doc:
                text += page.get_text()
        return text
    elif path.lower().endswith(".txt"):
        return open(path, encoding="utf-8", errors="ignore").read()
    return ""

def build_index_from_folder(root):
    index = {}
    for subdir, _, files in os.walk(root):
        for f in files:
            fp = os.path.join(subdir, f)
            try:
                index[fp] = extract_text_from_file(fp)
            except Exception:
                pass
    return index
