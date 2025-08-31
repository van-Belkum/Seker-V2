from typing import List, Optional
import io
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation

def extract_text_from_pdf_bytes(pdf_bytes: bytes) -> List[str]:
    pages = []
    with fitz.open(stream=pdf_bytes, filetype="pdf") as doc:
        for page in doc:
            pages.append(page.get_text("text"))
    return pages

def extract_text_from_pdf_path(path: str) -> List[str]:
    with open(path, "rb") as f:
        return extract_text_from_pdf_bytes(f.read())

def extract_text_from_docx_bytes(b: bytes) -> List[str]:
    mem = io.BytesIO(b)
    doc = Document(mem)
    text = []
    for p in doc.paragraphs:
        text.append(p.text)
    return ["\n".join(text)]

def extract_text_from_docx_path(path: str) -> List[str]:
    with open(path, "rb") as f:
        return extract_text_from_docx_bytes(f.read())

def extract_text_from_pptx_bytes(b: bytes) -> List[str]:
    mem = io.BytesIO(b)
    prs = Presentation(mem)
    pages = []
    for slide in prs.slides:
        t = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t.append(shape.text)
        pages.append("\n".join(t))
    return pages

def extract_text_from_pptx_path(path: str) -> List[str]:
    with open(path, "rb") as f:
        return extract_text_from_pptx_bytes(f.read())

def extract_text_from_txt_bytes(b: bytes) -> List[str]:
    try:
        s = b.decode("utf-8")
    except UnicodeDecodeError:
        s = b.decode("latin-1", errors="ignore")
    return [s]

def extract_text_from_file(name_or_path: str, raw_content: Optional[bytes]=None) -> List[str]:
    lower = name_or_path.lower()
    if lower.endswith(".pdf"):
        return extract_text_from_pdf_bytes(raw_content) if raw_content is not None else extract_text_from_pdf_path(name_or_path)
    if lower.endswith(".docx"):
        return extract_text_from_docx_bytes(raw_content) if raw_content is not None else extract_text_from_docx_path(name_or_path)
    if lower.endswith(".pptx"):
        return extract_text_from_pptx_bytes(raw_content) if raw_content is not None else extract_text_from_pptx_path(name_or_path)
    if lower.endswith(".txt"):
        return extract_text_from_txt_bytes(raw_content) if raw_content is not None else extract_text_from_txt_path(name_or_path)
    # fallback: treat as pdf
    return extract_text_from_pdf_bytes(raw_content) if raw_content is not None else extract_text_from_pdf_path(name_or_path)

def extract_text_from_txt_path(path: str) -> List[str]:
    with open(path, "rb") as f:
        return extract_text_from_txt_bytes(f.read())
