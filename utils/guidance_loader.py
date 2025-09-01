import io, zipfile, re, fitz, os
from docx import Document
from pptx import Presentation

KEYS = re.compile(r"\b(MUST|SHALL|REQUIRED|IMPORTANT NOTE|NOTE:)\b", re.I)

def _text_from_pdf_bytes(b: bytes) -> str:
    try:
        with fitz.open(stream=b, filetype="pdf") as doc:
            return "\n".join(page.get_text("text") for page in doc)
    except Exception:
        return ""

def _text_from_docx(b: bytes) -> str:
    f = io.BytesIO(b)
    doc = Document(f)
    return "\n".join(p.text for p in doc.paragraphs)

def _text_from_pptx(b: bytes) -> str:
    f = io.BytesIO(b)
    prs = Presentation(f)
    out = []
    for slide in prs.slides:
        for shp in slide.shapes:
            if hasattr(shp, "text"):
                out.append(shp.text)
    return "\n".join(out)

def lines_from_any(name: str, data: bytes) -> list[str]:
    n = name.lower()
    if n.endswith(".pdf"):
        txt = _text_from_pdf_bytes(data)
    elif n.endswith(".docx"):
        txt = _text_from_docx(data)
    elif n.endswith(".pptx"):
        txt = _text_from_pptx(data)
    else:
        return []
    # split to candidate rule lines
    lines = [l.strip() for l in txt.splitlines() if l.strip()]
    # keep informative-ish lines
    out = []
    for ln in lines:
        if len(ln) < 8: 
            continue
        if KEYS.search(ln) or any(k in ln.lower() for k in ["must","shall","required"]):
            out.append(ln)
    return out

def build_guidance_index_from_zip(zb: bytes) -> dict:
    z = zipfile.ZipFile(io.BytesIO(zb))
    rules = []
    guidance_hits = {}
    for zi in z.infolist():
        if zi.is_dir(): 
            continue
        if not zi.filename.lower().endswith((".pdf",".docx",".pptx")):
            continue
        data = z.read(zi)
        lines = lines_from_any(zi.filename, data)
        if not lines: 
            continue
        for ln in lines:
            rules.append({"source": zi.filename, "text": ln})
        guidance_hits[zi.filename] = len(lines)
    return {"rules": rules, "stats": guidance_hits}
